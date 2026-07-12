"""Document generation tool.

generate_document tool — builds Excel/Word/Markdown/slide deliverables
(via python-docx/openpyxl/python-pptx) and persists them through the
filesystem tool (FR-9).
"""

import base64
import io
import json
from typing import Literal

from langchain_core.tools import tool

from app.tools.filesystem_tool import write_file, write_binary_file, read_file, read_binary_file


@tool
async def generate_document(
    workspace_id: str,
    workspace_dir: str,
    path: str,
    format_type: Literal["excel", "word", "slide", "markdown"],
    content_data: str
) -> str:
    """Generate a document (Excel/Word/Markdown/Slide) and save it to the workspace.
    
    Args:
        workspace_id: The UUID of the workspace.
        workspace_dir: The host path of the workspace.
        path: The path to save the file (relative to the workspace).
        format_type: The format of the document (excel, word, slide, markdown).
        content_data: The content to write. For markdown, it's the raw text. 
                      For excel/word/slide, it must be a JSON string representing the data structure.
    """
    if format_type == "markdown":
        # Plain text
        return await write_file.invoke({"workspace_id": workspace_id, "workspace_dir": workspace_dir, "path": path, "content": content_data})
        
    try:
        data = json.loads(content_data)
    except json.JSONDecodeError as e:
        return f"Error: content_data must be a valid JSON string for {format_type} format. {str(e)}"
        
    buffer = io.BytesIO()
    
    try:
        if format_type == "excel":
            import openpyxl
            wb = openpyxl.Workbook()
            ws = wb.active
            if isinstance(data, list): # List of lists
                for row in data:
                    ws.append(row)
            wb.save(buffer)
            
        elif format_type == "word":
            import docx
            doc = docx.Document()
            if isinstance(data, list):
                for item in data:
                    if isinstance(item, dict) and "heading" in item:
                        doc.add_heading(item["heading"], level=item.get("level", 1))
                    elif isinstance(item, dict) and "paragraph" in item:
                        doc.add_paragraph(item["paragraph"])
                    elif isinstance(item, str):
                        doc.add_paragraph(item)
            else:
                doc.add_paragraph(str(data))
            doc.save(buffer)
            
        elif format_type == "slide":
            from pptx import Presentation
            prs = Presentation()
            if isinstance(data, list):
                for slide_data in data:
                    slide_layout = prs.slide_layouts[1] # title and content layout
                    slide = prs.slides.add_slide(slide_layout)
                    if isinstance(slide_data, dict):
                        if "title" in slide_data and slide.shapes.title:
                            slide.shapes.title.text = slide_data["title"]
                        if "content" in slide_data and len(slide.placeholders) > 1:
                            slide.placeholders[1].text = slide_data["content"]
            prs.save(buffer)
            
        else:
            return f"Error: Unsupported format type '{format_type}'."
            
        # ASSUMPTION: MCP tool arguments must be JSON-serializable, so binary data is passed as base64 strings.
        # Encode as base64 and write using MCP
        buffer.seek(0)
        base64_data = base64.b64encode(buffer.read()).decode("utf-8")
        
        return await write_binary_file.invoke({
            "workspace_id": workspace_id,
            "workspace_dir": workspace_dir,
            "path": path,
            "content_base64": base64_data
        })
        
    except ImportError as e:
        return f"Error: Missing required library for {format_type} generation: {str(e)}"
    except Exception as e:
        return f"Error generating document: {str(e)}"


@tool
async def edit_document(
    workspace_id: str,
    workspace_dir: str,
    path: str,
    format_type: Literal["word", "markdown"],
    operation: Literal["replace", "append"],
    target_text: str = "",
    replacement_text: str = ""
) -> str:
    """Edit an existing document in-place. Currently supports Markdown and Word files.
    
    Args:
        workspace_id: The UUID of the workspace.
        workspace_dir: The host path of the workspace.
        path: The path to the file to edit.
        format_type: The format of the document (markdown, word).
        operation: 'replace' to replace text, 'append' to add to the end.
        target_text: The exact text to find and replace (required if operation is 'replace').
        replacement_text: The new text to insert.
    """
    if format_type == "markdown":
        # Read file
        content = await read_file.invoke({"workspace_id": workspace_id, "workspace_dir": workspace_dir, "path": path})
        if content.startswith("Error:"):
            return content
            
        if operation == "replace":
            if target_text not in content:
                return f"Error: target_text not found in {path}"
            content = content.replace(target_text, replacement_text)
        elif operation == "append":
            content = content + "\n" + replacement_text
            
        return await write_file.invoke({
            "workspace_id": workspace_id,
            "workspace_dir": workspace_dir,
            "path": path,
            "content": content
        })
        
    elif format_type == "word":
        try:
            import docx
            
            # Read binary file via MCP
            base64_str = await read_binary_file.invoke({
                "workspace_id": workspace_id,
                "workspace_dir": workspace_dir,
                "path": path
            })
            if base64_str.startswith("Error:"):
                return base64_str
                
            binary_data = base64.b64decode(base64_str)
            buffer = io.BytesIO(binary_data)
            
            doc = docx.Document(buffer)
            
            if operation == "append":
                doc.add_paragraph(replacement_text)
            elif operation == "replace":
                replaced = False
                for para in doc.paragraphs:
                    if target_text in para.text:
                        para.text = para.text.replace(target_text, replacement_text)
                        replaced = True
                if not replaced:
                    return f"Error: target_text not found in Word document."
            
            # Save and write back
            out_buffer = io.BytesIO()
            doc.save(out_buffer)
            out_buffer.seek(0)
            
            out_base64 = base64.b64encode(out_buffer.read()).decode("utf-8")
            
            return await write_binary_file.invoke({
                "workspace_id": workspace_id,
                "workspace_dir": workspace_dir,
                "path": path,
                "content_base64": out_base64
            })
            
        except ImportError as e:
            return f"Error: Missing required library: {str(e)}"
        except Exception as e:
            return f"Error editing word document: {str(e)}"
            
    return f"Error: In-place edit for format '{format_type}' is not supported yet."
